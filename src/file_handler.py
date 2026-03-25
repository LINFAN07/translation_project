import io
import os
import fitz  # PyMuPDF
from docx import Document
from pdf_text_clean import clean_pdf_extracted_text

# 副檔名（小寫）集合，供 Streamlit file_uploader 與錯誤訊息使用
SUPPORTED_EXTENSIONS = frozenset(
    {".txt", ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt"}
)


def _read_txt_bytes(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _read_pdf_bytes(data: bytes) -> str:
    text = ""
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return clean_pdf_extracted_text(text)


def _read_docx_bytes(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    lines: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                lines.append("\t".join(cells))
    return "\n".join(lines)


def _read_xlsx_bytes(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    try:
        for sheet in wb.worksheets:
            lines.append(f"## {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(x.strip() for x in cells):
                    lines.append("\t".join(cells))
    finally:
        wb.close()
    return "\n".join(lines)


def _read_xls_bytes(data: bytes) -> str:
    import xlrd

    book = xlrd.open_workbook(file_contents=data)
    lines: list[str] = []
    for sheet in book.sheets():
        lines.append(f"## {sheet.name}")
        for r in range(sheet.nrows):
            row = sheet.row_values(r)
            cells = [str(c) for c in row]
            if any(c.strip() for c in cells):
                lines.append("\t".join(cells))
    return "\n".join(lines)


def _read_pptx_bytes(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"## 投影片 {i}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    lines.append(t)
    return "\n".join(lines)


def read_bytes_as_text(filename: str, data: bytes) -> str:
    """
    依檔名副檔名從二進位內容擷取純文字（供上傳檔與路徑讀檔共用）。
    """
    ext = os.path.splitext(filename)[1].lower()
    if not data:
        return ""

    if ext == ".txt":
        return _read_txt_bytes(data)
    if ext == ".pdf":
        return _read_pdf_bytes(data)
    if ext == ".docx":
        return _read_docx_bytes(data)
    if ext == ".doc":
        try:
            return _read_docx_bytes(data)
        except Exception as e:
            raise ValueError(
                "舊版 Word (.doc) 通常無法直接解析，請另存為 .docx 後再上傳。"
            ) from e
    if ext == ".xlsx":
        return _read_xlsx_bytes(data)
    if ext == ".xls":
        return _read_xls_bytes(data)
    if ext == ".pptx":
        return _read_pptx_bytes(data)
    if ext == ".ppt":
        raise ValueError(
            "舊版 PowerPoint (.ppt) 不支援，請在 PowerPoint 中另存為 .pptx 後再上傳。"
        )

    raise ValueError(
        f"不支援的檔案格式: {ext}。支援: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
    )


def read_file_content(file_path: str) -> str:
    """從檔案路徑讀取並轉成純文字。"""
    name = os.path.basename(file_path)
    with open(file_path, "rb") as f:
        data = f.read()
    return read_bytes_as_text(name, data)


def read_streamlit_upload(uploaded_file) -> str:
    """Streamlit UploadedFile：以檔名與位元組擷取文字。"""
    return read_bytes_as_text(uploaded_file.name, uploaded_file.getvalue())
